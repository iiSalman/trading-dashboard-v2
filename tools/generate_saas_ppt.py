#!/usr/bin/env python3
"""Generate SAAS_PROPOSAL.pptx — 14-slide pivot proposal deck.

Re-run any time to regenerate. Uses the dashboard's color tokens so the deck
visually matches the live product (https://iisalman.github.io/trading-dashboard-v2).

Single dep: python-pptx (see tools/requirements-ppt.txt).
"""
from __future__ import annotations

from datetime import date
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# ── dashboard palette ─────────────────────────────────────────────
BG       = RGBColor(0x0B, 0x0F, 0x1A)
BG2      = RGBColor(0x0F, 0x14, 0x24)
SURFACE  = RGBColor(0x11, 0x18, 0x27)
SURFACE2 = RGBColor(0x1A, 0x22, 0x35)
BORDER   = RGBColor(0x1F, 0x2D, 0x42)
TEXT     = RGBColor(0xE2, 0xE8, 0xF0)
TEXT_DIM = RGBColor(0x94, 0xA3, 0xB8)
MUTED    = RGBColor(0x64, 0x74, 0x8B)
GREEN    = RGBColor(0x10, 0xB9, 0x81)
RED      = RGBColor(0xEF, 0x44, 0x44)
YELLOW   = RGBColor(0xF5, 0x9E, 0x0B)
BLUE     = RGBColor(0x3B, 0x82, 0xF6)
PURPLE   = RGBColor(0x8B, 0x5C, 0xF6)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

FONT_MONO = "Menlo"     # macOS native mono
FONT_SANS = "Helvetica Neue"

OUTPUT_PATH = Path(__file__).resolve().parent.parent / "SAAS_PROPOSAL.pptx"
TOTAL_SLIDES = 14


# ── primitives ────────────────────────────────────────────────────
def blank_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG
    return slide


def add_text(slide, x, y, w, h, text, *,
             size=14, bold=False, color=TEXT, align=PP_ALIGN.LEFT,
             font=FONT_SANS, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor
    lines = text.split("\n") if "\n" in text else [text]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.name = font
        run.font.color.rgb = color
    return box


def add_rect(slide, x, y, w, h, *, fill=SURFACE, border=BORDER, border_w=0.75,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    s = slide.shapes.add_shape(shape, x, y, w, h)
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        s.adjustments[0] = 0.04
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if border is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = border
        s.line.width = Pt(border_w)
    s.shadow.inherit = False
    return s


def add_arrow(slide, x1, y1, x2, y2, color=BLUE, weight=2.5):
    line = slide.shapes.add_connector(2, x1, y1, x2, y2)  # 2 = straight
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    return line


def add_logo_chip(slide, x, y, size_in=0.32):
    # blue "Q" chip matching dashboard
    chip = add_rect(slide, x, y, Inches(size_in), Inches(size_in),
                    fill=BLUE, border=None, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text(slide, x, y, Inches(size_in), Inches(size_in), "Q",
             size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
             font=FONT_MONO, anchor=MSO_ANCHOR.MIDDLE)


def header_strip(slide, slide_num, title_kicker, title_main):
    """Top bar: Q logo + small mono kicker + bold title + slide-counter top-right."""
    add_logo_chip(slide, Inches(0.45), Inches(0.32))
    add_text(slide, Inches(0.90), Inches(0.30), Inches(2.5), Inches(0.25),
             "QUANT DASHBOARD", size=10, color=TEXT_DIM, font=FONT_MONO, bold=True)
    add_text(slide, Inches(0.90), Inches(0.50), Inches(8.0), Inches(0.25),
             title_kicker, size=9, color=MUTED, font=FONT_MONO)
    # bottom thin divider
    add_rect(slide, Inches(0.45), Inches(0.95), Inches(12.45), Inches(0.01),
             fill=BORDER, border=None, shape=MSO_SHAPE.RECTANGLE)
    # slide counter top-right
    add_text(slide, Inches(11.5), Inches(0.32), Inches(1.5), Inches(0.30),
             f"{slide_num:02d} / {TOTAL_SLIDES:02d}", size=10, color=MUTED,
             font=FONT_MONO, align=PP_ALIGN.RIGHT)
    # big title under the divider
    add_text(slide, Inches(0.45), Inches(1.10), Inches(12.45), Inches(0.65),
             title_main, size=30, bold=True, color=TEXT, font=FONT_SANS)


def footer(slide):
    add_text(slide, Inches(0.45), Inches(7.10), Inches(12.45), Inches(0.30),
             "quant-dashboard · saas pivot proposal · " + date.today().isoformat(),
             size=9, color=MUTED, font=FONT_MONO)


# ── slides ────────────────────────────────────────────────────────
def slide_01_cover(prs):
    s = blank_slide(prs)
    # accent strip on the left
    add_rect(s, Inches(0), Inches(0), Inches(0.25), SLIDE_H,
             fill=BLUE, border=None, shape=MSO_SHAPE.RECTANGLE)
    add_logo_chip(s, Inches(0.80), Inches(0.80), size_in=0.55)
    add_text(s, Inches(1.55), Inches(0.78), Inches(6), Inches(0.35),
             "QUANT DASHBOARD", size=14, bold=True, color=TEXT_DIM, font=FONT_MONO)
    add_text(s, Inches(1.55), Inches(1.04), Inches(6), Inches(0.30),
             "options analytics · live GEX · score-driven picks",
             size=11, color=MUTED, font=FONT_MONO)

    # hero title
    add_text(s, Inches(0.80), Inches(2.40), Inches(11.5), Inches(1.10),
             "From shared snapshot", size=54, bold=True, color=TEXT, font=FONT_SANS)
    add_text(s, Inches(0.80), Inches(3.30), Inches(11.5), Inches(1.10),
             "to live multi-tenant SaaS.", size=54, bold=True, color=BLUE, font=FONT_SANS)

    add_text(s, Inches(0.80), Inches(4.70), Inches(11.5), Inches(0.50),
             "A pivot blueprint for the trading-dashboard-v2 project.",
             size=18, color=TEXT_DIM, font=FONT_SANS)

    # bottom meta strip
    add_rect(s, Inches(0.80), Inches(6.20), Inches(11.5), Inches(0.55),
             fill=SURFACE, border=BORDER)
    add_text(s, Inches(1.05), Inches(6.30), Inches(8), Inches(0.35),
             "prepared for  iiSalman",
             size=12, color=TEXT, font=FONT_MONO, bold=True,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(8.0), Inches(6.30), Inches(4.3), Inches(0.35),
             date.today().isoformat() + "  ·  v1.0",
             size=12, color=TEXT_DIM, font=FONT_MONO, align=PP_ALIGN.RIGHT,
             anchor=MSO_ANCHOR.MIDDLE)


def slide_02_pitch(prs):
    s = blank_slide(prs)
    header_strip(s, 2, "executive summary", "The pitch, in four bullets.")

    bullets = [
        ("01", "Target market",
         "U.S. retail/semi-pro options traders, 0DTE-curious, PDT-aware.",
         BLUE),
        ("02", "Differentiator",
         "GEX-first, score-driven, live. No clutter, no upsell soup.",
         GREEN),
        ("03", "Ask",
         "$1k seed budget + 8 focused weeks. Solo bootstrappable.",
         YELLOW),
        ("04", "Milestone",
         "Paid beta in 8 weeks. Break-even at 13 users on a $29/mo plan.",
         PURPLE),
    ]
    y = Inches(1.95)
    for num, label, body, color in bullets:
        card = add_rect(s, Inches(0.45), y, Inches(12.45), Inches(1.10),
                        fill=SURFACE, border=BORDER)
        # color chip with num
        add_rect(s, Inches(0.70), y + Inches(0.20), Inches(0.70), Inches(0.70),
                 fill=color, border=None)
        add_text(s, Inches(0.70), y + Inches(0.20), Inches(0.70), Inches(0.70),
                 num, size=20, bold=True, color=BG, font=FONT_MONO,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(1.65), y + Inches(0.15), Inches(10.5), Inches(0.35),
                 label.upper(), size=11, bold=True, color=color, font=FONT_MONO)
        add_text(s, Inches(1.65), y + Inches(0.50), Inches(10.5), Inches(0.55),
                 body, size=17, color=TEXT, font=FONT_SANS)
        y += Inches(1.25)

    footer(s)


def slide_03_today(prs):
    s = blank_slide(prs)
    header_strip(s, 3, "current state", "Where we are today.")

    # architecture diagram: 5 boxes left-to-right
    boxes = [
        ("CRON",            "GitHub Actions\n*/5 * * * *",          BLUE),
        ("DATA",            "yfinance\n(Yahoo scrape)",             YELLOW),
        ("BUILD",           "build_snapshot.py\n→ snapshot.json",   GREEN),
        ("HOST",            "GitHub Pages\nstatic site",            PURPLE),
        ("USER",            "Browser\nfetches JSON",                TEXT_DIM),
    ]
    box_w = Inches(2.20)
    box_h = Inches(1.40)
    gap = Inches(0.20)
    total = box_w * 5 + gap * 4
    start_x = (SLIDE_W - total) / 2
    y = Inches(2.30)
    cx = []
    for i, (lab, body, color) in enumerate(boxes):
        x = start_x + (box_w + gap) * i
        add_rect(s, x, y, box_w, box_h, fill=SURFACE, border=color, border_w=1.5)
        add_text(s, x, y + Inches(0.10), box_w, Inches(0.30),
                 lab, size=11, bold=True, color=color, font=FONT_MONO,
                 align=PP_ALIGN.CENTER)
        add_text(s, x, y + Inches(0.45), box_w, box_h - Inches(0.50),
                 body, size=13, color=TEXT, font=FONT_MONO,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        cx.append(x + box_w)
    # arrows between
    for i in range(4):
        x1 = cx[i]
        x2 = cx[i] + gap
        add_arrow(s, x1, y + box_h / 2, x2, y + box_h / 2, color=TEXT_DIM, weight=1.5)

    # commentary panel below
    panel_y = Inches(4.30)
    add_rect(s, Inches(0.45), panel_y, Inches(12.45), Inches(2.30),
             fill=BG2, border=BORDER)
    add_text(s, Inches(0.75), panel_y + Inches(0.20), Inches(11.5), Inches(0.30),
             "WHAT IT IS", size=11, bold=True, color=TEXT_DIM, font=FONT_MONO)
    body = (
        "A free, single-page dashboard. One shared JSON snapshot, rebuilt every "
        "5 minutes by a cron job that scrapes Yahoo Finance and commits the "
        "result back to git. Anyone with the URL sees the same data."
    )
    add_text(s, Inches(0.75), panel_y + Inches(0.55), Inches(11.5), Inches(1.0),
             body, size=15, color=TEXT, font=FONT_SANS)
    add_text(s, Inches(0.75), panel_y + Inches(1.55), Inches(11.5), Inches(0.30),
             "GREAT FOR", size=11, bold=True, color=GREEN, font=FONT_MONO)
    add_text(s, Inches(0.75), panel_y + Inches(1.85), Inches(11.5), Inches(0.40),
             "personal use, demos, hobby projects.",
             size=15, color=TEXT, font=FONT_SANS)
    footer(s)


def slide_04_blockers(prs):
    s = blank_slide(prs)
    header_strip(s, 4, "the problem", "Why it can't be a SaaS as-is.")

    rows = [
        ("Yahoo ToS violation",
         "yfinance scrapes Yahoo. Yahoo's terms forbid commercial redistribution. Day-one legal risk.",
         RED),
        ("Shared snapshot",
         "Every visitor sees the same data at the same staleness. Paying customers expect per-user views.",
         YELLOW),
        ("No authentication",
         "No login, no accounts, no per-user state. Can't bill what we can't identify.",
         YELLOW),
        ("No payments",
         "No Stripe, no subscription tiers, no entitlement gating. Free is the only price.",
         YELLOW),
        ("No streaming",
         "Architecture commits JSON to git on a 5-min cron. No path to push live updates to a browser.",
         RED),
    ]
    y = Inches(1.90)
    row_h = Inches(0.95)
    for i, (label, body, color) in enumerate(rows):
        bg = SURFACE if i % 2 == 0 else BG2
        add_rect(s, Inches(0.45), y, Inches(12.45), row_h, fill=bg, border=BORDER)
        # severity pill
        add_rect(s, Inches(0.70), y + Inches(0.27), Inches(0.45), Inches(0.40),
                 fill=color, border=None)
        sev = "✕" if color == RED else "!"
        add_text(s, Inches(0.70), y + Inches(0.27), Inches(0.45), Inches(0.40),
                 sev, size=16, bold=True, color=BG, font=FONT_MONO,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(1.30), y + Inches(0.12), Inches(4.0), Inches(0.40),
                 label, size=15, bold=True, color=color, font=FONT_SANS,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(5.30), y + Inches(0.12), Inches(7.5), row_h - Inches(0.24),
                 body, size=14, color=TEXT, font=FONT_SANS,
                 anchor=MSO_ANCHOR.MIDDLE)
        y += row_h + Inches(0.05)
    footer(s)


def slide_05_customer(prs):
    s = blank_slide(prs)
    header_strip(s, 5, "the customer", "Who pays, and why.")

    # left: persona card
    add_rect(s, Inches(0.45), Inches(1.95), Inches(5.50), Inches(4.90),
             fill=SURFACE, border=BORDER)
    add_text(s, Inches(0.75), Inches(2.15), Inches(5.0), Inches(0.30),
             "PERSONA", size=11, bold=True, color=BLUE, font=FONT_MONO)
    add_text(s, Inches(0.75), Inches(2.45), Inches(5.0), Inches(0.50),
             "Semi-pro retail trader", size=22, bold=True, color=TEXT, font=FONT_SANS)
    add_text(s, Inches(0.75), Inches(3.05), Inches(5.0), Inches(0.40),
             "(25–45, U.S., $50k–500k account)", size=12, color=TEXT_DIM, font=FONT_MONO)

    persona_lines = [
        "•  trades U.S. options 2–10x / week",
        "•  reads /r/options, watches Tasty Live",
        "•  pays for Unusual Whales OR Whalestream",
        "•  already comfortable reading Greeks",
        "•  obsessive about GEX & 0DTE flow",
        "•  hates noise, wants signal",
    ]
    add_text(s, Inches(0.75), Inches(3.55), Inches(5.0), Inches(3.0),
             "\n".join(persona_lines),
             size=14, color=TEXT, font=FONT_SANS)

    # right: competitive landscape
    add_rect(s, Inches(6.20), Inches(1.95), Inches(6.70), Inches(4.90),
             fill=BG2, border=BORDER)
    add_text(s, Inches(6.50), Inches(2.15), Inches(6.0), Inches(0.30),
             "WHAT THEY PAY FOR TODAY", size=11, bold=True, color=GREEN, font=FONT_MONO)
    competitors = [
        ("Unusual Whales",    "$50–80 / mo",   "options flow + sweeps"),
        ("Whalestream",       "$60–150 / mo",  "dark pool + flow"),
        ("OptionStrat",       "$30–60 / mo",   "strategy builder"),
        ("SpotGamma",         "$99–349 / mo",  "GEX & vanna research"),
        ("Tasty Live",        "$0 (broker)",   "education + watchlists"),
    ]
    y = Inches(2.55)
    for name, price, blurb in competitors:
        add_text(s, Inches(6.50), y, Inches(2.8), Inches(0.35),
                 name, size=14, bold=True, color=TEXT, font=FONT_SANS)
        add_text(s, Inches(9.30), y, Inches(1.6), Inches(0.35),
                 price, size=13, color=YELLOW, font=FONT_MONO)
        add_text(s, Inches(10.90), y, Inches(2.0), Inches(0.35),
                 blurb, size=12, color=TEXT_DIM, font=FONT_SANS)
        y += Inches(0.50)

    # our wedge
    add_rect(s, Inches(6.20), Inches(5.30), Inches(6.70), Inches(1.55),
             fill=SURFACE, border=GREEN, border_w=1.5)
    add_text(s, Inches(6.50), Inches(5.45), Inches(6.0), Inches(0.30),
             "OUR WEDGE", size=11, bold=True, color=GREEN, font=FONT_MONO)
    add_text(s, Inches(6.50), Inches(5.75), Inches(6.0), Inches(1.0),
             "GEX-first scoring. One number per ticker tells you whether to fade or follow. No 12 panels of clutter. $29 starting price undercuts SpotGamma 10×.",
             size=13, color=TEXT, font=FONT_SANS)
    footer(s)


def slide_06_architecture(prs):
    s = blank_slide(prs)
    header_strip(s, 6, "target architecture", "Live data, multi-tenant, streaming.")

    # data source (top)
    add_rect(s, Inches(0.45), Inches(1.95), Inches(3.5), Inches(1.00),
             fill=SURFACE, border=YELLOW, border_w=1.5)
    add_text(s, Inches(0.45), Inches(2.05), Inches(3.5), Inches(0.30),
             "DATA PROVIDER", size=10, bold=True, color=YELLOW, font=FONT_MONO,
             align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.45), Inches(2.35), Inches(3.5), Inches(0.55),
             "Theta Data\nWebSocket · options chain",
             size=13, color=TEXT, font=FONT_MONO, align=PP_ALIGN.CENTER)

    # backend (middle)
    add_rect(s, Inches(4.50), Inches(1.95), Inches(4.50), Inches(2.40),
             fill=SURFACE, border=BLUE, border_w=2.0)
    add_text(s, Inches(4.50), Inches(2.05), Inches(4.50), Inches(0.30),
             "BACKEND  ·  Python FastAPI on Railway", size=10, bold=True,
             color=BLUE, font=FONT_MONO, align=PP_ALIGN.CENTER)
    backend_lines = [
        "▸ subscribe to Theta WebSocket",
        "▸ compute Δ Γ PCR GEX score (analyzer.py)",
        "▸ cache hot state in Upstash Redis",
        "▸ broadcast via Server-Sent Events",
        "▸ entitlement check per request",
    ]
    add_text(s, Inches(4.75), Inches(2.45), Inches(4.10), Inches(1.85),
             "\n".join(backend_lines), size=12, color=TEXT, font=FONT_MONO)

    # frontend (right)
    add_rect(s, Inches(9.55), Inches(1.95), Inches(3.40), Inches(2.40),
             fill=SURFACE, border=GREEN, border_w=2.0)
    add_text(s, Inches(9.55), Inches(2.05), Inches(3.40), Inches(0.30),
             "FRONTEND  ·  Next.js on Vercel", size=10, bold=True,
             color=GREEN, font=FONT_MONO, align=PP_ALIGN.CENTER)
    front_lines = [
        "▸ Next.js 15 + Tailwind",
        "▸ EventSource subscriber",
        "▸ per-user watchlist",
        "▸ tier-gated UI",
        "▸ Stripe Checkout / Portal",
    ]
    add_text(s, Inches(9.80), Inches(2.45), Inches(3.10), Inches(1.85),
             "\n".join(front_lines), size=12, color=TEXT, font=FONT_MONO)

    # arrows
    add_arrow(s, Inches(3.95), Inches(2.45), Inches(4.50), Inches(2.45),
              color=YELLOW, weight=2.0)
    add_arrow(s, Inches(9.00), Inches(3.15), Inches(9.55), Inches(3.15),
              color=BLUE, weight=2.0)

    # side rails bottom: Auth, DB, Billing
    rail_y = Inches(4.85)
    rails = [
        ("AUTH",      "Supabase Auth\nemail · OAuth · MAU billing",  PURPLE),
        ("DATABASE",  "Supabase Postgres\nusers · subs · watchlists", PURPLE),
        ("BILLING",   "Stripe Checkout\nsubs · invoices · webhooks",  PURPLE),
        ("HOSTING",   "Railway + Vercel\n+ Upstash + GH Pages docs",  PURPLE),
    ]
    rail_w = Inches(3.0)
    for i, (label, body, color) in enumerate(rails):
        x = Inches(0.45) + (rail_w + Inches(0.12)) * i
        add_rect(s, x, rail_y, rail_w, Inches(1.40),
                 fill=BG2, border=BORDER)
        add_text(s, x, rail_y + Inches(0.12), rail_w, Inches(0.30),
                 label, size=10, bold=True, color=color, font=FONT_MONO,
                 align=PP_ALIGN.CENTER)
        add_text(s, x, rail_y + Inches(0.45), rail_w, Inches(0.85),
                 body, size=12, color=TEXT, font=FONT_MONO,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)

    # callout: latency budget
    add_text(s, Inches(0.45), Inches(6.45), Inches(12.45), Inches(0.40),
             "End-to-end latency budget: ~1.5s from CBOE → Theta → backend → user screen.",
             size=12, color=TEXT_DIM, font=FONT_MONO, align=PP_ALIGN.CENTER)
    footer(s)


def slide_07_data_providers(prs):
    s = blank_slide(prs)
    header_strip(s, 7, "decision · data provider",
                 "Who's allowed to redistribute to paying customers?")

    rows = [
        ("Theta Data",         "~$250",   "✓", "✓", "✓ retail SaaS", "★ recommended",        GREEN),
        ("Polygon Massive",    "~$500+",  "✓", "✓", "✓ B2B",        "sales-gated; pricey",   YELLOW),
        ("CBOE LiveVol",       "sales",   "✓", "✓", "✓ enterprise", "institutional",         YELLOW),
        ("Polygon Developer",  "$199",    "✓", "✓", "✕ internal only", "no redistribution",  RED),
        ("Tradier",            "$0*",     "✓", "✓", "✕ per-user",   "broker account each",   RED),
        ("yfinance / Yahoo",   "$0",      "△", "✕", "✕ ToS forbids","day-one legal risk",    RED),
    ]
    headers = ["Provider", "Monthly", "Real-time", "Stream", "Redistribution", "Verdict"]
    col_x = [Inches(0.45), Inches(2.90), Inches(4.60), Inches(6.10), Inches(7.60), Inches(10.10)]
    col_w = [Inches(2.45), Inches(1.70), Inches(1.50), Inches(1.50), Inches(2.50), Inches(2.80)]

    y = Inches(2.00)
    # header row
    add_rect(s, Inches(0.45), y, Inches(12.45), Inches(0.45),
             fill=SURFACE2, border=BORDER)
    for h, x, w in zip(headers, col_x, col_w):
        add_text(s, x, y, w, Inches(0.45), h.upper(),
                 size=11, bold=True, color=TEXT_DIM, font=FONT_MONO,
                 anchor=MSO_ANCHOR.MIDDLE)
    y += Inches(0.45)

    for i, row in enumerate(rows):
        bg = BG2 if i % 2 == 0 else SURFACE
        add_rect(s, Inches(0.45), y, Inches(12.45), Inches(0.55),
                 fill=bg, border=BORDER)
        provider, monthly, rt, stream, redist, verdict, color = row
        cells = [provider, monthly, rt, stream, redist, verdict]
        for cell, x, w in zip(cells, col_x, col_w):
            font_color = color if cell == verdict else TEXT
            font = FONT_SANS if cell in (provider, verdict, redist) else FONT_MONO
            bold = (cell == provider)
            add_text(s, x, y, w, Inches(0.55), cell,
                     size=13, color=font_color, font=font, bold=bold,
                     anchor=MSO_ANCHOR.MIDDLE)
        y += Inches(0.55)

    # callout
    add_rect(s, Inches(0.45), Inches(6.10), Inches(12.45), Inches(0.85),
             fill=SURFACE, border=GREEN, border_w=1.5)
    add_text(s, Inches(0.75), Inches(6.20), Inches(11.5), Inches(0.30),
             "RECOMMENDATION", size=11, bold=True, color=GREEN, font=FONT_MONO)
    add_text(s, Inches(0.75), Inches(6.50), Inches(11.5), Inches(0.45),
             "Start on Theta Data Standard ($250/mo). Confirm redistribution scope in writing before launch.",
             size=14, color=TEXT, font=FONT_SANS)
    footer(s)


def slide_08_stack(prs):
    s = blank_slide(prs)
    header_strip(s, 8, "tech stack", "Boring stack, fast to ship.")

    columns = [
        ("BACKEND", BLUE, [
            ("Python 3.12",    "reuse existing analyzer.py"),
            ("FastAPI",        "async REST + SSE endpoints"),
            ("websockets",     "Theta Data subscriber loop"),
            ("Redis (Upstash)", "hot per-ticker state"),
            ("Postgres",       "users + audit + history"),
        ]),
        ("FRONTEND", GREEN, [
            ("Next.js 15",     "App Router, RSC"),
            ("React 19",       "UI"),
            ("Tailwind",       "match dashboard tokens"),
            ("TanStack Query", "client cache + retry"),
            ("EventSource",    "subscribe to backend SSE"),
        ]),
        ("PLATFORM", PURPLE, [
            ("Railway",        "$20/mo backend host"),
            ("Vercel",         "free → $20 frontend host"),
            ("Supabase",       "$25/mo Auth + Postgres"),
            ("Upstash Redis",  "$10/mo serverless cache"),
            ("Stripe",         "subs billing, 2.9% + 30¢"),
        ]),
    ]
    col_w = Inches(4.05)
    gap = Inches(0.15)
    total = col_w * 3 + gap * 2
    start_x = (SLIDE_W - total) / 2
    y = Inches(2.00)
    for i, (label, color, items) in enumerate(columns):
        x = start_x + (col_w + gap) * i
        add_rect(s, x, y, col_w, Inches(4.85), fill=SURFACE, border=BORDER)
        add_rect(s, x, y, col_w, Inches(0.55), fill=color, border=None)
        add_text(s, x, y, col_w, Inches(0.55), label,
                 size=14, bold=True, color=BG, font=FONT_MONO,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        yy = y + Inches(0.70)
        for name, desc in items:
            add_text(s, x + Inches(0.20), yy, col_w - Inches(0.40), Inches(0.30),
                     name, size=14, bold=True, color=TEXT, font=FONT_MONO)
            add_text(s, x + Inches(0.20), yy + Inches(0.30), col_w - Inches(0.40), Inches(0.40),
                     desc, size=12, color=TEXT_DIM, font=FONT_SANS)
            yy += Inches(0.78)
    footer(s)


def slide_09_pricing(prs):
    s = blank_slide(prs)
    header_strip(s, 9, "pricing", "Three tiers, one decision per user.")

    tiers = [
        ("FREE", "$0", "always free", MUTED, BG2, [
            "15-min delayed data",
            "top 5 tickers only",
            "basic GEX & PCR",
            "no watchlist save",
            "no alerts",
        ], "lead gen"),
        ("BASIC", "$29", "/ month", BLUE, SURFACE, [
            "real-time options chain",
            "25-ticker watchlist",
            "full score table",
            "GEX glossary + tools",
            "email support",
        ], "the wedge"),
        ("PRO", "$79", "/ month", GREEN, SURFACE, [
            "everything in Basic",
            "custom watchlists",
            "price + flow alerts",
            "historical backtest",
            "REST API access",
        ], "power user · margin"),
    ]
    col_w = Inches(4.05)
    gap = Inches(0.15)
    total = col_w * 3 + gap * 2
    start_x = (SLIDE_W - total) / 2
    y = Inches(1.90)
    for i, (name, price, period, color, bg, features, sub) in enumerate(tiers):
        x = start_x + (col_w + gap) * i
        # tier card
        border_w = 2.0 if name == "BASIC" else 1.0
        add_rect(s, x, y, col_w, Inches(4.95), fill=bg, border=color, border_w=border_w)
        # ribbon for recommended
        if name == "BASIC":
            add_rect(s, x + col_w - Inches(1.5), y - Inches(0.05), Inches(1.5), Inches(0.35),
                     fill=color, border=None)
            add_text(s, x + col_w - Inches(1.5), y - Inches(0.05), Inches(1.5), Inches(0.35),
                     "RECOMMENDED", size=9, bold=True, color=BG, font=FONT_MONO,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # header
        add_text(s, x + Inches(0.20), y + Inches(0.35), col_w - Inches(0.40), Inches(0.40),
                 name, size=13, bold=True, color=color, font=FONT_MONO)
        add_text(s, x + Inches(0.20), y + Inches(0.80), col_w - Inches(0.40), Inches(0.75),
                 price, size=44, bold=True, color=TEXT, font=FONT_SANS)
        add_text(s, x + Inches(0.20), y + Inches(1.65), col_w - Inches(0.40), Inches(0.30),
                 period, size=12, color=TEXT_DIM, font=FONT_MONO)
        add_text(s, x + Inches(0.20), y + Inches(1.95), col_w - Inches(0.40), Inches(0.30),
                 "→  " + sub, size=11, color=color, font=FONT_MONO)
        # divider
        add_rect(s, x + Inches(0.20), y + Inches(2.35), col_w - Inches(0.40), Inches(0.02),
                 fill=BORDER, border=None, shape=MSO_SHAPE.RECTANGLE)
        # features
        yy = y + Inches(2.55)
        for f in features:
            add_text(s, x + Inches(0.30), yy, col_w - Inches(0.50), Inches(0.30),
                     "·  " + f, size=12, color=TEXT, font=FONT_SANS)
            yy += Inches(0.38)
    footer(s)


def slide_10_costs(prs):
    s = blank_slide(prs)
    header_strip(s, 10, "cost model", "Fixed burn, simple break-even.")

    # left: fixed monthly table
    add_rect(s, Inches(0.45), Inches(1.95), Inches(6.0), Inches(4.90),
             fill=SURFACE, border=BORDER)
    add_text(s, Inches(0.75), Inches(2.10), Inches(5.5), Inches(0.30),
             "FIXED MONTHLY COSTS", size=11, bold=True, color=TEXT_DIM, font=FONT_MONO)

    costs = [
        ("Theta Data Standard",      "$250"),
        ("Supabase Pro (Auth + DB)", "$25"),
        ("Railway (backend)",        "$20"),
        ("Vercel Pro (frontend)",    "$20"),
        ("Upstash Redis",            "$10"),
        ("Domain + email",           "$3"),
    ]
    yy = Inches(2.55)
    for label, amt in costs:
        add_text(s, Inches(0.75), yy, Inches(4.5), Inches(0.40),
                 label, size=14, color=TEXT, font=FONT_SANS, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(5.25), yy, Inches(1.0), Inches(0.40),
                 amt, size=14, bold=True, color=YELLOW, font=FONT_MONO,
                 align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
        yy += Inches(0.55)
    # divider + total
    add_rect(s, Inches(0.75), yy + Inches(0.05), Inches(5.5), Inches(0.02),
             fill=BORDER, border=None, shape=MSO_SHAPE.RECTANGLE)
    yy += Inches(0.20)
    add_text(s, Inches(0.75), yy, Inches(4.5), Inches(0.45),
             "TOTAL", size=15, bold=True, color=TEXT, font=FONT_MONO,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(5.25), yy, Inches(1.0), Inches(0.45),
             "$328", size=20, bold=True, color=GREEN, font=FONT_MONO,
             align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(0.75), yy + Inches(0.55), Inches(5.5), Inches(0.35),
             "+ Stripe 2.9 % + 30¢ per charge (variable)",
             size=11, color=TEXT_DIM, font=FONT_MONO)

    # right: break-even table
    add_rect(s, Inches(6.70), Inches(1.95), Inches(6.20), Inches(4.90),
             fill=BG2, border=BORDER)
    add_text(s, Inches(7.00), Inches(2.10), Inches(5.5), Inches(0.30),
             "BREAK-EVEN  ·  users to cover fixed burn", size=11, bold=True,
             color=GREEN, font=FONT_MONO)

    headers = ["Plan", "Price", "Users", "@ 50 users", "@ 200 users"]
    col_x = [Inches(7.00), Inches(8.30), Inches(9.40), Inches(10.50), Inches(11.70)]
    col_w = [Inches(1.30), Inches(1.10), Inches(1.10), Inches(1.20), Inches(1.20)]

    yy = Inches(2.55)
    for h, x, w in zip(headers, col_x, col_w):
        add_text(s, x, yy, w, Inches(0.35), h.upper(),
                 size=10, bold=True, color=TEXT_DIM, font=FONT_MONO,
                 anchor=MSO_ANCHOR.MIDDLE)
    yy += Inches(0.40)

    rows = [
        ("Basic",   "$29",   "12",  "$1,122",  "$5,472"),
        ("Mid",     "$49",   "7",   "$2,122",  "$9,472"),
        ("Pro",     "$79",   "5",   "$3,622",  "$15,472"),
    ]
    colors_row = [BLUE, YELLOW, GREEN]
    for (plan, price, users, m50, m200), color in zip(rows, colors_row):
        add_rect(s, Inches(7.00), yy, Inches(5.9), Inches(0.55),
                 fill=SURFACE, border=BORDER)
        cells = [plan, price, users, m50, m200]
        for cell, x, w in zip(cells, col_x, col_w):
            font_color = color if cell == plan else TEXT
            add_text(s, x, yy, w, Inches(0.55), cell,
                     size=13, bold=(cell == plan), color=font_color,
                     font=FONT_MONO, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        yy += Inches(0.60)

    # callout
    add_rect(s, Inches(7.00), Inches(5.40), Inches(5.9), Inches(1.30),
             fill=SURFACE, border=GREEN, border_w=1.5)
    add_text(s, Inches(7.20), Inches(5.55), Inches(5.5), Inches(0.30),
             "TAKEAWAY", size=10, bold=True, color=GREEN, font=FONT_MONO)
    add_text(s, Inches(7.20), Inches(5.85), Inches(5.5), Inches(0.80),
             "12 paying users covers everything. 50 users = ~$1k/mo MRR. The pure data cost ($250) dominates everything else.",
             size=12, color=TEXT, font=FONT_SANS)
    footer(s)


def slide_11_rollout(prs):
    s = blank_slide(prs)
    header_strip(s, 11, "phased rollout", "Eight weeks to paid beta. Six months to full product.")

    phases = [
        ("PHASE 1  ·  Weeks 1–6",
         "MVP launch — single tier",
         BLUE,
         [
             "sign Theta Data contract (verify redistribution scope)",
             "scaffold FastAPI service polling Theta REST every 30s",
             "Next.js shell + Supabase auth + Stripe single $29 plan",
             "migrate analyzer.py compute into the FastAPI worker",
             "private beta with 5–10 users from existing audience",
         ]),
        ("PHASE 2  ·  Weeks 7–12",
         "Real-time + Pro tier",
         GREEN,
         [
             "swap REST polling for Theta WebSocket subscriber",
             "push live updates to browser via SSE / WebSocket",
             "ship Pro $79 tier with custom watchlists",
             "Stripe customer portal for self-serve upgrades",
             "public launch on Twitter / r/options / IndieHackers",
         ]),
        ("PHASE 3  ·  Months 4–6",
         "Retention features",
         PURPLE,
         [
             "price + flow alerts (SMS via Twilio, Discord webhook)",
             "historical backtest replay (TimescaleDB-backed)",
             "REST + WebSocket API for developers (Pro tier)",
             "team plans / multi-seat",
             "case study + content marketing flywheel",
         ]),
    ]

    y = Inches(1.95)
    h = Inches(1.55)
    for kicker, title, color, items in phases:
        # left header band
        add_rect(s, Inches(0.45), y, Inches(3.20), h, fill=SURFACE, border=color, border_w=1.5)
        add_text(s, Inches(0.65), y + Inches(0.20), Inches(2.9), Inches(0.40),
                 kicker, size=11, bold=True, color=color, font=FONT_MONO)
        add_text(s, Inches(0.65), y + Inches(0.55), Inches(2.9), Inches(0.50),
                 title, size=20, bold=True, color=TEXT, font=FONT_SANS)
        # right bullet list
        add_rect(s, Inches(3.80), y, Inches(9.10), h, fill=BG2, border=BORDER)
        bullets_txt = "\n".join(["▸  " + it for it in items])
        add_text(s, Inches(4.00), y + Inches(0.20), Inches(8.8), h - Inches(0.30),
                 bullets_txt, size=13, color=TEXT, font=FONT_SANS)
        y += h + Inches(0.10)
    footer(s)


def slide_12_risks(prs):
    s = blank_slide(prs)
    header_strip(s, 12, "risks & mitigations", "What kills this, and how to not die.")

    risks = [
        ("Data licensing scope",
         "Theta's $250 tier may not cover retail SaaS redistribution at scale.",
         "Get scope confirmed in writing from sales before launch. Budget for $500/mo tier.",
         RED),
        ("Regulatory framing",
         "'BULLISH / BEARISH' labels could be construed as investment advice.",
         "Soften copy. Add strong disclaimers. ~$1k for a securities lawyer review pre-launch.",
         RED),
        ("Yahoo cliff",
         "Current dashboard uses yfinance — illegal once we charge for the same data.",
         "Must migrate to Theta before flipping the paywall on. No grace period.",
         RED),
        ("WebSocket scaling",
         "At 100+ concurrent users the FastAPI process needs horizontal scaling.",
         "Use Railway autoscale + a fanout layer (Cloudflare Workers / Pusher) when MRR > $1k.",
         YELLOW),
        ("Beta retention",
         "Users may try it, then churn after 30 days. Common in fintech SaaS.",
         "Build alerts + watchlist sync (Phase 2) early — these drive habit & lock-in.",
         YELLOW),
    ]

    y = Inches(1.95)
    h = Inches(0.95)
    for risk, problem, mitigation, color in risks:
        add_rect(s, Inches(0.45), y, Inches(12.45), h, fill=SURFACE, border=BORDER)
        # severity bar
        add_rect(s, Inches(0.45), y, Inches(0.18), h, fill=color, border=None,
                 shape=MSO_SHAPE.RECTANGLE)
        add_text(s, Inches(0.75), y + Inches(0.12), Inches(3.2), Inches(0.35),
                 risk, size=14, bold=True, color=color, font=FONT_SANS)
        add_text(s, Inches(0.75), y + Inches(0.45), Inches(3.2), h - Inches(0.50),
                 problem, size=11, color=TEXT_DIM, font=FONT_SANS)
        # mitigation panel
        add_rect(s, Inches(4.10), y + Inches(0.10), Inches(8.70), h - Inches(0.20),
                 fill=BG2, border=None)
        add_text(s, Inches(4.30), y + Inches(0.18), Inches(8.40), Inches(0.30),
                 "MITIGATION", size=9, bold=True, color=GREEN, font=FONT_MONO)
        add_text(s, Inches(4.30), y + Inches(0.45), Inches(8.40), h - Inches(0.55),
                 mitigation, size=13, color=TEXT, font=FONT_SANS)
        y += h + Inches(0.08)
    footer(s)


def slide_13_actions(prs):
    s = blank_slide(prs)
    header_strip(s, 13, "next steps", "Five things to do this week.")

    actions = [
        ("Email Theta Data sales",
         "Confirm $250 Standard tier explicitly permits retail SaaS redistribution. Get it in writing.",
         "30 min"),
        ("Stand up a private repo",
         "github.com/iiSalman/dashboard-saas — separate from the public dashboard.",
         "10 min"),
        ("Scaffold FastAPI service",
         "/health, /tickers, /quote/{symbol}. REST polling stub against yfinance to validate the loop.",
         "3–4 hr"),
        ("Spin up Next.js + Supabase",
         "create-next-app + Supabase project + email auth + protected route. Stripe test mode wired.",
         "4–6 hr"),
        ("Register a domain",
         "Something punchy and memorable: gexedge.com, quantpip.io, flowflag.app — under $15 / yr.",
         "20 min"),
    ]

    y = Inches(2.00)
    h = Inches(0.85)
    for i, (label, body, eff) in enumerate(actions, start=1):
        add_rect(s, Inches(0.45), y, Inches(12.45), h, fill=SURFACE, border=BORDER)
        # checkbox circle
        add_rect(s, Inches(0.70), y + Inches(0.20), Inches(0.45), Inches(0.45),
                 fill=BG2, border=BLUE, border_w=1.5, shape=MSO_SHAPE.OVAL)
        add_text(s, Inches(0.70), y + Inches(0.20), Inches(0.45), Inches(0.45),
                 str(i), size=15, bold=True, color=BLUE, font=FONT_MONO,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(1.35), y + Inches(0.14), Inches(8.0), Inches(0.35),
                 label, size=14, bold=True, color=TEXT, font=FONT_SANS)
        add_text(s, Inches(1.35), y + Inches(0.45), Inches(8.0), Inches(0.40),
                 body, size=12, color=TEXT_DIM, font=FONT_SANS)
        # effort pill
        add_rect(s, Inches(11.20), y + Inches(0.27), Inches(1.45), Inches(0.35),
                 fill=BG2, border=GREEN, border_w=1.0)
        add_text(s, Inches(11.20), y + Inches(0.27), Inches(1.45), Inches(0.35),
                 eff, size=11, color=GREEN, font=FONT_MONO,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        y += h + Inches(0.10)
    footer(s)


def slide_14_sources(prs):
    s = blank_slide(prs)
    header_strip(s, 14, "appendix", "Sources & further reading.")

    sources = [
        ("DATA PROVIDERS", [
            ("Theta Data — pricing",         "thetadata.net/pricing"),
            ("Theta Data — commercial use",  "thetadata.net/commercial-use"),
            ("Polygon Business / Massive",   "polygon.io/business"),
            ("CBOE All Access API",          "datashop.cboe.com/cboe-all-access-api"),
        ]),
        ("STACK", [
            ("FastAPI",                       "fastapi.tiangolo.com"),
            ("Supabase",                      "supabase.com / pricing"),
            ("Stripe Subscriptions",          "stripe.com/docs/billing/subscriptions"),
            ("Railway, Vercel, Upstash",      "railway.app · vercel.com · upstash.com"),
        ]),
        ("LEGAL & MARKET", [
            ("Yahoo Finance ToS",             "policies.yahoo.com/us/en/yahoo/terms"),
            ("SEC — investment-adviser guidance",
                                              "sec.gov/investment/adviser-guidance"),
            ("SpotGamma blog (GEX research)",  "spotgamma.com/blog"),
            ("Unusual Whales (competitor)",    "unusualwhales.com"),
        ]),
    ]

    col_w = Inches(4.05)
    gap = Inches(0.15)
    total = col_w * 3 + gap * 2
    start_x = (SLIDE_W - total) / 2
    y = Inches(2.00)
    for i, (label, items) in enumerate(sources):
        x = start_x + (col_w + gap) * i
        add_rect(s, x, y, col_w, Inches(4.40), fill=SURFACE, border=BORDER)
        add_text(s, x, y + Inches(0.20), col_w, Inches(0.30),
                 label, size=11, bold=True, color=BLUE, font=FONT_MONO,
                 align=PP_ALIGN.CENTER)
        yy = y + Inches(0.65)
        for name, url in items:
            add_text(s, x + Inches(0.25), yy, col_w - Inches(0.50), Inches(0.30),
                     name, size=13, bold=True, color=TEXT, font=FONT_SANS)
            add_text(s, x + Inches(0.25), yy + Inches(0.32), col_w - Inches(0.50), Inches(0.30),
                     url, size=10, color=TEXT_DIM, font=FONT_MONO)
            yy += Inches(0.85)

    # closing line
    add_text(s, Inches(0.45), Inches(6.65), Inches(12.45), Inches(0.40),
             "questions? iterate on this deck — re-run tools/generate_saas_ppt.py to regenerate.",
             size=11, color=MUTED, font=FONT_MONO, align=PP_ALIGN.CENTER)
    footer(s)


# ── main ──────────────────────────────────────────────────────────
def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_01_cover(prs)
    slide_02_pitch(prs)
    slide_03_today(prs)
    slide_04_blockers(prs)
    slide_05_customer(prs)
    slide_06_architecture(prs)
    slide_07_data_providers(prs)
    slide_08_stack(prs)
    slide_09_pricing(prs)
    slide_10_costs(prs)
    slide_11_rollout(prs)
    slide_12_risks(prs)
    slide_13_actions(prs)
    slide_14_sources(prs)

    prs.save(OUTPUT_PATH)
    print(str(OUTPUT_PATH))


if __name__ == "__main__":
    main()
